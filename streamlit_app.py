import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import io
import datetime
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, BubbleChartData
from lxml import etree
from sklearn.linear_model import LogisticRegression
from sklearn.model_selection import train_test_split
from sklearn.metrics import classification_report
import config as cfg

BLUE = "#29B5E8"
DARK = "#262730"
RED = "#D94032"
GREEN = "#2ECC71"
AMBER = "#FFBE2E"
LIGHT_GRAY = "#F0F2F6"
MID_GRAY = "#888888"
BG_CARD = "#1B1F2A"

st.set_page_config(
    page_title=cfg.APP_TITLE,
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
    [data-testid="stMetric"] {
        background: linear-gradient(135deg, #1B1F2A 0%, #262730 100%);
        border: 1px solid #333;
        border-radius: 12px;
        padding: 16px 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.3);
    }
    [data-testid="stMetricLabel"] {
        font-size: 0.85rem !important;
        color: #AAA !important;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    [data-testid="stMetricValue"] {
        font-size: 1.8rem !important;
        font-weight: 700 !important;
    }
    div[data-testid="stHorizontalBlock"] > div {
        padding: 0 4px;
    }
    .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #0E1117 0%, #1B1F2A 100%);
    }
    h1, h2, h3 {
        font-weight: 700 !important;
    }
    .stDownloadButton > button {
        background: linear-gradient(135deg, #29B5E8 0%, #1a8ab5 100%) !important;
        color: white !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 12px 32px !important;
        font-weight: 600 !important;
        font-size: 1rem !important;
        transition: all 0.2s ease !important;
        box-shadow: 0 4px 12px rgba(41,181,232,0.3) !important;
    }
    .stDownloadButton > button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 6px 16px rgba(41,181,232,0.4) !important;
    }
    div.stExpander {
        border: 1px solid #333;
        border-radius: 12px;
        background: #1B1F2A;
    }
</style>
""", unsafe_allow_html=True)

conn = st.connection("snowflake", ttl=os.getenv("SNOWFLAKE_CONNECTION_TTL"))
session = conn.session()


@st.cache_data(ttl=600)
def load_summary():
    return session.sql(f"SELECT * FROM {cfg.SUMMARY_TABLE}").to_pandas()


@st.cache_data(ttl=600)
def load_collections():
    return session.sql(f"""
        SELECT DELINQUENCY_BUCKET, TOTAL_ACCOUNTS, TOTAL_OUTSTANDING,
               AVG_DPD_DAYS, COLLECTION_RATE_PCT, CURE_RATE_PCT, ROLL_RATE_PCT
        FROM {cfg.COLLECTIONS_TABLE}
        ORDER BY AVG_DPD_DAYS
    """).to_pandas()


@st.cache_data(ttl=600)
def load_performance():
    return session.sql(f"""
        SELECT ORIGINATION_MONTH, RISK_TIER, LOAN_COUNT, TOTAL_FUNDED,
               DELINQUENT_LOANS, DELINQUENCY_RATE_PCT, AVG_APR, AVG_CREDIT_SCORE, CNL_RATE_PCT
        FROM {cfg.PERFORMANCE_TABLE}
        ORDER BY ORIGINATION_MONTH
    """).to_pandas()


@st.cache_data(ttl=600)
def load_risk_tier_scatter():
    return session.sql(f"""
        SELECT RISK_TIER,
               ROUND(AVG(AVG_CREDIT_SCORE), 0) AS AVG_CREDIT_SCORE,
               ROUND(AVG(DELINQUENCY_RATE_PCT), 2) AS AVG_DELINQUENCY_RATE,
               SUM(TOTAL_FUNDED) AS TOTAL_FUNDED
        FROM {cfg.PERFORMANCE_TABLE}
        GROUP BY RISK_TIER
    """).to_pandas()


@st.cache_data(ttl=600)
def run_cortex_summary():
    summary = session.sql(f"SELECT * FROM {cfg.SUMMARY_TABLE}").to_pandas()
    tiers = session.sql(f"""
        SELECT RISK_TIER, COUNT(*) AS MONTHS, ROUND(AVG(DELINQUENCY_RATE_PCT),2) AS AVG_DELINQ,
               ROUND(AVG(CNL_RATE_PCT),2) AS AVG_CNL, ROUND(AVG(AVG_CREDIT_SCORE),0) AS AVG_SCORE
        FROM {cfg.PERFORMANCE_TABLE} GROUP BY RISK_TIER
    """).to_pandas()
    prompt = (
        "You are a senior credit risk analyst. Summarize this auto loan portfolio. "
        "Highlight key risks, trends, and actions. Use bullet points. Be concise.\n\n"
        f"Portfolio: {summary.to_dict(orient='records')}\n\n"
        f"By tier: {tiers.to_dict(orient='records')}"
    )
    prompt_escaped = prompt.replace("'", "''")
    result = session.sql(f"SELECT SNOWFLAKE.CORTEX.COMPLETE('{cfg.CORTEX_MODEL}', '{prompt_escaped}') AS SUMMARY").to_pandas()
    return result["SUMMARY"].iloc[0]


def train_risk_model(perf):
    loans = perf.dropna(subset=["AVG_CREDIT_SCORE", "AVG_APR", "DELINQUENCY_RATE_PCT"])
    if len(loans) < 10:
        return None, None, ""
    X = loans[["AVG_CREDIT_SCORE", "AVG_APR"]].astype(float).values
    y = (loans["DELINQUENCY_RATE_PCT"] > 5).astype(int).values
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.3, random_state=42)
    model = LogisticRegression(random_state=42)
    model.fit(X_train, y_train)
    y_pred = model.predict(X_test)
    report = classification_report(y_test, y_pred, target_names=["Not Delinquent", "Delinquent"])
    return model, loans, report


with st.sidebar:
    st.markdown(f"""
    <div style="text-align:center; padding: 1.5rem 0;">
        <h1 style="color:{BLUE}; margin:0; font-size:1.8rem;">❄️ {cfg.APP_TITLE}</h1>
        <p style="color:{MID_GRAY}; margin:4px 0 0 0; font-size:0.85rem; letter-spacing:1px;">{cfg.APP_SUBTITLE}</p>
    </div>
    """, unsafe_allow_html=True)
    st.divider()

    if st.button("🔄 Refresh Data", use_container_width=True):
        st.cache_data.clear()
        st.rerun()

    st.divider()
    perf = load_performance()
    all_tiers = sorted(perf["RISK_TIER"].dropna().unique().tolist())
    selected_tiers = st.multiselect("Risk Tiers", all_tiers, default=all_tiers)

    st.divider()
    st.markdown(f"""
    <div style="padding:0.5rem 0; font-size:0.75rem; color:{MID_GRAY};">
        <p>Powered by Snowflake<br>Dynamic Tables + Cortex AI</p>
        <p>{datetime.date.today().strftime('%B %d, %Y')}</p>
    </div>
    """, unsafe_allow_html=True)

summary = load_summary()
collections = load_collections()
perf_filtered = perf[perf["RISK_TIER"].isin(selected_tiers)] if selected_tiers else perf

st.markdown(f"# 📊 {cfg.APP_TITLE}")
st.markdown(f"<p style='color:{MID_GRAY}; margin-top:-10px;'>Real-time portfolio monitoring powered by Snowflake Dynamic Tables</p>", unsafe_allow_html=True)

row = summary.iloc[0]
total_funded_m = float(row["TOTAL_FUNDED_AMOUNT"]) / 1_000_000
outstanding_m = float(row["TOTAL_OUTSTANDING_BALANCE"]) / 1_000_000
avg_delinquency = round(float(perf_filtered["DELINQUENCY_RATE_PCT"].mean()), 1)

c1, c2, c3, c4 = st.columns(4)
c1.metric("Total Loans", f"{int(row['TOTAL_LOANS']):,}")
c2.metric("Total Funded", f"${total_funded_m:.1f}M")
c3.metric("Outstanding", f"${outstanding_m:.1f}M")
c4.metric("Avg Credit Score", f"{int(row['WEIGHTED_AVG_CREDIT_SCORE'])}")

c5, c6, c7, c8 = st.columns(4)
c5.metric("Current %", f"{row['CURRENT_PCT']}%")
c6.metric("30+ DPD", f"{row['DPD_30_PCT']}%", delta=f"{row['DPD_30_PCT']}%", delta_color="inverse")
c7.metric("Avg Delinquency", f"{avg_delinquency}%", delta=f"{avg_delinquency}%", delta_color="inverse")
c8.metric("Projected Losses", f"${float(row['PROJECTED_LOSSES']):,.0f}", delta="at risk", delta_color="inverse")

st.markdown("")

col_left, col_right = st.columns(2)

with col_left:
    st.subheader("Delinquency Distribution")
    bucket_order = ["CURRENT", "30_DPD", "60_DPD", "90_DPD", "120+_DPD"]
    coll_sorted = collections.copy()
    coll_sorted["sort_key"] = coll_sorted["DELINQUENCY_BUCKET"].map(
        {b: i for i, b in enumerate(bucket_order)}
    )
    coll_sorted = coll_sorted.sort_values("sort_key")
    fig_collections = go.Figure()
    fig_collections.add_trace(go.Bar(
        x=coll_sorted["DELINQUENCY_BUCKET"],
        y=coll_sorted["TOTAL_ACCOUNTS"],
        marker_color=[GREEN if b == "CURRENT" else AMBER if "30" in b else RED
                       for b in coll_sorted["DELINQUENCY_BUCKET"]],
        text=coll_sorted["TOTAL_ACCOUNTS"],
        textposition="outside",
        textfont=dict(size=12),
    ))
    fig_collections.update_layout(
        xaxis_title="Delinquency Bucket",
        yaxis_title="Total Accounts",
        height=420,
        margin=dict(t=20, b=40),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(color="#FAFAFA"),
        xaxis=dict(gridcolor="#333"),
        yaxis=dict(gridcolor="#333"),
    )
    st.plotly_chart(fig_collections, use_container_width=True)

with col_right:
    st.subheader("Vintage Delinquency Curves")
    tier_colors = {"PRIME": BLUE, "NEAR_PRIME": AMBER, "SUBPRIME": RED,
                   "SUPER_PRIME": GREEN, "DEEP_SUBPRIME": "#FF6B6B"}
    fig_vintage = go.Figure()
    for tier in sorted(perf_filtered["RISK_TIER"].unique()):
        tier_data = perf_filtered[perf_filtered["RISK_TIER"] == tier].sort_values("ORIGINATION_MONTH")
        fig_vintage.add_trace(go.Scatter(
            x=tier_data["ORIGINATION_MONTH"],
            y=tier_data["DELINQUENCY_RATE_PCT"],
            mode="lines+markers",
            name=tier,
            line=dict(color=tier_colors.get(tier, MID_GRAY), width=2),
            marker=dict(size=5),
        ))
    fig_vintage.update_layout(
        xaxis_title="Origination Month",
        yaxis_title="Delinquency Rate (%)",
        height=420,
        margin=dict(t=20, b=40),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(color="#FAFAFA"),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        xaxis=dict(gridcolor="#333"),
        yaxis=dict(gridcolor="#333"),
    )
    st.plotly_chart(fig_vintage, use_container_width=True)

col_left2, col_right2 = st.columns(2)

with col_left2:
    st.subheader("Risk Tier: Score vs Delinquency")
    scatter_data = load_risk_tier_scatter()
    if selected_tiers:
        scatter_data = scatter_data[scatter_data["RISK_TIER"].isin(selected_tiers)]
    fig_scatter = px.scatter(
        scatter_data,
        x="AVG_CREDIT_SCORE",
        y="AVG_DELINQUENCY_RATE",
        size="TOTAL_FUNDED",
        color="RISK_TIER",
        color_discrete_map=tier_colors,
        hover_data=["TOTAL_FUNDED"],
        height=420,
    )
    fig_scatter.update_layout(
        xaxis_title="Avg Credit Score",
        yaxis_title="Avg Delinquency Rate (%)",
        margin=dict(t=20, b=40),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(color="#FAFAFA"),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        xaxis=dict(gridcolor="#333"),
        yaxis=dict(gridcolor="#333"),
    )
    st.plotly_chart(fig_scatter, use_container_width=True)

with col_right2:
    st.subheader("Delinquency Probability Model")
    model, model_data, report = train_risk_model(perf_filtered)
    if model is not None:
        scores = np.linspace(550, 850, 100)
        probs = model.predict_proba(np.column_stack([scores, np.full(100, 10.0)]))[:, 1]
        fig_model = go.Figure()
        fig_model.add_trace(go.Scatter(
            x=scores, y=probs * 100,
            mode="lines",
            line=dict(color=RED, width=3),
            fill="tozeroy",
            fillcolor="rgba(217,64,50,0.15)",
        ))
        fig_model.update_layout(
            xaxis_title="Credit Score",
            yaxis_title="Delinquency Probability (%)",
            height=420,
            margin=dict(t=20, b=40),
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)",
            font=dict(color="#FAFAFA"),
            xaxis=dict(gridcolor="#333"),
            yaxis=dict(gridcolor="#333"),
        )
        st.plotly_chart(fig_model, use_container_width=True)
    else:
        st.info("Insufficient data for risk model with selected filters.")
        fig_model = None

st.divider()

with st.expander("🤖 AI Executive Summary — Snowflake Cortex", expanded=False):
    if st.button("Generate AI Summary", key="ai_btn"):
        with st.status(f"Calling Snowflake Cortex ({cfg.CORTEX_MODEL})...", expanded=True) as status:
            try:
                st.write("Fetching portfolio data...")
                ai_text = run_cortex_summary()
                st.session_state["ai_text"] = ai_text
                status.update(label="AI Summary complete", state="complete")
            except Exception as e:
                st.session_state["ai_text"] = f"Error: {e}"
                status.update(label="AI Summary failed", state="error")
    if "ai_text" in st.session_state:
        st.markdown(st.session_state["ai_text"])

st.divider()


def build_pptx(summary_row, collections_df, perf_df, scatter_df, model_obj, ai_summary):
    PBLUE = RGBColor(0x29, 0xB5, 0xE8)
    PDARK = RGBColor(0x26, 0x27, 0x30)
    PWHITE = RGBColor(0xFF, 0xFF, 0xFF)
    PLGRAY = RGBColor(0xF0, 0xF2, 0xF6)
    PRED = RGBColor(0xD9, 0x40, 0x32)
    PGREEN = RGBColor(0x2E, 0xCC, 0x71)
    PAMBER = RGBColor(0xFF, 0xBE, 0x2E)
    PMGRAY = RGBColor(0x88, 0x88, 0x88)

    TIER_PPTX_COLORS = {
        "PRIME": PBLUE, "NEAR_PRIME": PAMBER, "SUBPRIME": PRED,
        "SUPER_PRIME": PGREEN, "DEEP_SUBPRIME": RGBColor(0xFF, 0x6B, 0x6B),
    }

    nsmap = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
             "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    def _add_rect(slide, l, t, w, h, rgb):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.line.fill.background()
        s.fill.solid()
        s.fill.fore_color.rgb = rgb
        return s

    def _add_text(slide, text, l, t, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        txb.word_wrap = wrap
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txb

    def _slide_header(prs, header):
        blank = prs.slide_layouts[6]
        s = prs.slides.add_slide(blank)
        _add_rect(s, 0, 0, SW, 0.55, PDARK)
        _add_text(s, header, 0.2, 0.05, SW - 0.4, 0.45, 15, bold=True, color=PWHITE)
        return s, None

    def _set_point_color(series_element, idx, rgb):
        dPt = etree.SubElement(series_element, etree.QName(nsmap["c"], "dPt"))
        idx_el = etree.SubElement(dPt, etree.QName(nsmap["c"], "idx"))
        idx_el.set("val", str(idx))
        spPr = etree.SubElement(dPt, etree.QName(nsmap["c"], "spPr"))
        solidFill = etree.SubElement(spPr, etree.QName(nsmap["a"], "solidFill"))
        srgbClr = etree.SubElement(solidFill, etree.QName(nsmap["a"], "srgbClr"))
        srgbClr.set("val", str(rgb))

    def _set_series_color(series, rgb):
        spPr = series._element.find(etree.QName(nsmap["c"], "spPr"))
        if spPr is None:
            spPr = etree.SubElement(series._element, etree.QName(nsmap["c"], "spPr"))
        for old in spPr.findall(etree.QName(nsmap["a"], "solidFill")):
            spPr.remove(old)
        solidFill = etree.SubElement(spPr, etree.QName(nsmap["a"], "solidFill"))
        srgbClr = etree.SubElement(solidFill, etree.QName(nsmap["a"], "srgbClr"))
        srgbClr.set("val", str(rgb))
        ln = spPr.find(etree.QName(nsmap["a"], "ln"))
        if ln is None:
            ln = etree.SubElement(spPr, etree.QName(nsmap["a"], "ln"))
        for old in ln.findall(etree.QName(nsmap["a"], "solidFill")):
            ln.remove(old)
        ln_fill = etree.SubElement(ln, etree.QName(nsmap["a"], "solidFill"))
        ln_clr = etree.SubElement(ln_fill, etree.QName(nsmap["a"], "srgbClr"))
        ln_clr.set("val", str(rgb))

    r = summary_row
    total_funded_m = float(r["TOTAL_FUNDED_AMOUNT"]) / 1_000_000
    avg_del = round(float(perf_df["DELINQUENCY_RATE_PCT"].mean()), 1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW = 13.333

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, 7.5, PBLUE)
    _add_rect(s, 0, 5.8, SW, 1.7, PDARK)
    _add_text(s, cfg.APP_TITLE, 0.7, 1.6, SW - 1.4, 1.1, 40, bold=True, color=PWHITE)
    _add_text(s, "Exploratory Data Analysis \u2014 Snowflake Streamlit", 0.7, 2.75, SW - 1.4, 0.7, 22, color=PWHITE)
    _add_text(s, datetime.date.today().strftime("%B %Y"), 0.7, 3.55, 4, 0.5, 16, color=RGBColor(0xCC, 0xEE, 0xF8))
    _add_text(s, "Powered by Snowflake Dynamic Tables + Cortex AI", 0.7, 6.15, SW - 1.4, 0.5, 13, color=PMGRAY)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, 0.55, PDARK)
    _add_text(s, "Portfolio KPIs", 0.2, 0.05, SW - 0.4, 0.45, 15, bold=True, color=PWHITE)
    kpis_r1 = [
        ("Total loans", f"{int(r['TOTAL_LOANS']):,}", PBLUE),
        ("Total funded", f"${total_funded_m:.2f}M", PBLUE),
        ("Outstanding", f"${float(r['TOTAL_OUTSTANDING_BALANCE'])/1e6:.2f}M", PBLUE),
        ("Avg credit score", f"{int(r['WEIGHTED_AVG_CREDIT_SCORE'])}", PBLUE),
    ]
    kpis_r2 = [
        ("Current %", f"{r['CURRENT_PCT']}%", PGREEN),
        ("30+ DPD", f"{r['DPD_30_PCT']}%", PAMBER),
        ("Avg delinquency", f"{avg_del}%", PRED),
        ("Projected losses", f"${float(r['PROJECTED_LOSSES']):,.0f}", PRED),
    ]
    box_w = (SW - 0.36 - 3 * 0.13) / 4
    box_gap = 0.13
    for row_kpis, box_top in [(kpis_r1, 0.72), (kpis_r2, 1.85)]:
        for i, (label, value, val_color) in enumerate(row_kpis):
            lx = 0.18 + i * (box_w + box_gap)
            _add_rect(s, lx, box_top, box_w, 0.95, PLGRAY)
            _add_text(s, value, lx+0.1, box_top+0.04, box_w-0.2, 0.5, 20, bold=True, color=val_color, align=PP_ALIGN.CENTER)
            _add_text(s, label, lx+0.1, box_top+0.56, box_w-0.2, 0.3, 9, color=PDARK, align=PP_ALIGN.CENTER)

    bucket_order = ["CURRENT", "30_DPD", "60_DPD", "90_DPD", "120+_DPD"]
    cdf = collections_df.copy()
    cdf["sort_key"] = cdf["DELINQUENCY_BUCKET"].map({b: i for i, b in enumerate(bucket_order)})
    cdf = cdf.sort_values("sort_key")
    bucket_colors = []
    for b in cdf["DELINQUENCY_BUCKET"]:
        if b == "CURRENT":
            bucket_colors.append(PGREEN)
        elif "30" in b:
            bucket_colors.append(PAMBER)
        else:
            bucket_colors.append(PRED)

    chart_data = CategoryChartData()
    chart_data.categories = list(cdf["DELINQUENCY_BUCKET"])
    chart_data.add_series("Total Accounts", list(cdf["TOTAL_ACCOUNTS"].astype(float)))

    s3, _ = _slide_header(prs, "Delinquency Distribution by Bucket")
    chart_frame = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.18), Inches(0.65), Inches(SW - 0.36), Inches(6.6), chart_data)
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series_el = plot.series[0]._element
    for idx, clr in enumerate(bucket_colors):
        _set_point_color(series_el, idx, clr)
    plot.series[0].data_labels.show_value = True
    plot.series[0].data_labels.font.size = Pt(9)
    plot.series[0].data_labels.font.bold = True
    plot.series[0].data_labels.number_format = "#,##0"

    tiers = sorted(perf_df["RISK_TIER"].dropna().unique())
    months = sorted(perf_df["ORIGINATION_MONTH"].dropna().unique())
    month_labels = [str(m)[:10] if hasattr(m, 'strftime') else str(m) for m in months]

    line_data = CategoryChartData()
    line_data.categories = month_labels
    for tier in tiers:
        td = perf_df[perf_df["RISK_TIER"] == tier].copy()
        td_grouped = td.groupby("ORIGINATION_MONTH")["DELINQUENCY_RATE_PCT"].mean()
        vals = []
        for m in months:
            if m in td_grouped.index:
                vals.append(float(td_grouped[m]))
            else:
                vals.append(None)
        line_data.add_series(tier, vals)

    s4, _ = _slide_header(prs, "Delinquency Rate by Vintage Month and Risk Tier")
    chart_frame = s4.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.18), Inches(0.65), Inches(SW - 0.36), Inches(6.6), line_data)
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    for i, tier in enumerate(tiers):
        ser = chart.series[i]
        clr = TIER_PPTX_COLORS.get(tier, PMGRAY)
        _set_series_color(ser, clr)
        ser.smooth = False
        ser.format.line.width = Pt(2)

    bubble_data = BubbleChartData()
    for _, brow in scatter_df.iterrows():
        tier_name = brow["RISK_TIER"]
        ser = bubble_data.add_series(tier_name)
        funded = float(brow["TOTAL_FUNDED"])
        ser.add_data_point(float(brow["AVG_CREDIT_SCORE"]), float(brow["AVG_DELINQUENCY_RATE"]), funded)

    s5, _ = _slide_header(prs, "Risk Tier: Credit Score vs Delinquency Rate")
    chart_frame = s5.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE, Inches(0.18), Inches(0.65), Inches(SW - 0.36), Inches(6.6), bubble_data)
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    for i, (_, brow) in enumerate(scatter_df.iterrows()):
        clr = TIER_PPTX_COLORS.get(brow["RISK_TIER"], PMGRAY)
        _set_series_color(chart.series[i], clr)

    if model_obj is not None:
        scores_arr = np.linspace(550, 850, 30)
        probs_arr = model_obj.predict_proba(np.column_stack([scores_arr, np.full(30, 10.0)]))[:, 1] * 100
        prob_data = CategoryChartData()
        prob_data.categories = [str(int(s)) for s in scores_arr]
        prob_data.add_series("Delinquency Probability %", [round(float(p), 2) for p in probs_arr])

        s6, _ = _slide_header(prs, "Predicted Delinquency Probability by Credit Score")
        chart_frame = s6.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(0.18), Inches(0.65), Inches(SW - 0.36), Inches(6.6), prob_data)
        chart = chart_frame.chart
        chart.has_legend = False
        _set_series_color(chart.series[0], PRED)
        chart.series[0].smooth = True
        chart.series[0].format.line.width = Pt(3)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, SW, 0.55, PDARK)
    _add_text(s, "AI-Generated Executive Summary (Snowflake Cortex)", 0.2, 0.05, SW - 0.4, 0.45, 15, bold=True, color=PWHITE)
    _add_text(s, ai_summary, 0.3, 0.75, SW - 0.6, 6.5, 11, color=PDARK, wrap=True)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


st.subheader("📥 Export to PowerPoint")
st.markdown(f"<p style='color:{MID_GRAY};'>Generate a presentation deck with all KPIs, charts, and AI summary.</p>", unsafe_allow_html=True)

if st.button("⚡ Generate PowerPoint", use_container_width=False):
    with st.spinner("Building presentation..."):
        ai_summary = st.session_state.get("ai_text", "")
        if not ai_summary:
            try:
                ai_summary = run_cortex_summary()
            except Exception:
                ai_summary = "AI summary unavailable."

        pptx_bytes = build_pptx(
            row, collections, perf_filtered, load_risk_tier_scatter(),
            model, ai_summary,
        )
        st.session_state["pptx_bytes"] = pptx_bytes
        st.session_state["pptx_ready"] = True

if st.session_state.get("pptx_ready"):
    st.download_button(
        label="⬇️  Download Presentation",
        data=st.session_state["pptx_bytes"],
        file_name=f"{cfg.DOWNLOAD_PREFIX}_{datetime.date.today().strftime('%Y%m%d')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    st.success(f"Presentation ready — click above to download.")
